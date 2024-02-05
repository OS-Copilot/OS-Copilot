from fastapi import APIRouter, Depends
from pydantic import BaseModel
from pptx import Presentation
import os
import time
import requests

router = APIRouter()

CWD = os.getcwd()  # path of current working directory
LIB_DIR = os.path.dirname(__file__)  # path of library
TEMPLATE_DIR = os.path.join(LIB_DIR, "templates")  # path of templates
CACHE_DIR = os.path.join(CWD, "cache")  # path of cache_dir
IMAGE_BED_PATTERN = 'https://source.unsplash.com/featured/?{}'  # url pattern for image bed

if not os.path.exists(CACHE_DIR):
    os.makedirs(CACHE_DIR)

ppt_file = None  # a pointer to the powerpoint object


class CreateFileModel(BaseModel):
    theme: str


class GetImageModel(BaseModel):
    keywords: str


class AddFirstPageModel(BaseModel):
    title: str
    subtitle: str


class AddTextPageModel(BaseModel):
    title: str
    bullet_items: str


class AddTextImagePageModel(BaseModel):
    title: str
    bullet_items: str
    image: str


@router.post("/tools/ppt/create_file")
async def create_file(item: CreateFileModel):
    global ppt_file
    ppt_file = Presentation(os.path.join(TEMPLATE_DIR, f"{item.theme}.pptx"))
    return "created a ppt file."


@router.post("/tools/ppt/get_image")
async def get_image(item: GetImageModel):
    picture_url = IMAGE_BED_PATTERN.format(item.keywords)
    response = requests.get(picture_url)
    img_local_path = os.path.join(CACHE_DIR, f"{time.time()}.jpg")
    with open(img_local_path, 'wb') as f:
        f.write(response.content)
    return img_local_path


@router.post("/tools/ppt/add_first_page")
async def add_first_page(item: AddFirstPageModel):
    global ppt_file
    slide = ppt_file.slides.add_slide(ppt_file.slide_layouts[0])  # layout for first page (title and subtitle only)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    title_shape.text = item.title
    subtitle_shape.text = item.subtitle
    return "added first page."


@router.post("/tools/ppt/add_text_page")
async def add_text_page(item: AddTextPageModel):
    global ppt_file
    slide = ppt_file.slides.add_slide(ppt_file.slide_layouts[1])
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    title_shape.text = item.title
    tf = body_shape.text_frame
    bullet_items = item.bullet_items.split("[SPAN]")
    for bullet_item in bullet_items:
        bullet_item_strip = bullet_item.strip()
        p = tf.add_paragraph()
        p.text = bullet_item_strip
        p.level = 1
    return "added text page."


@router.post("/tools/ppt/add_text_image_page")
async def add_text_image_page(item: AddTextImagePageModel):
    global ppt_file
    slide = ppt_file.slides.add_slide(ppt_file.slide_layouts[3])
    title_shape = slide.shapes.title
    title_shape.text = item.title
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    bullet_items = item.bullet_items.split("[SPAN]")
    for bullet_item in bullet_items:
        bullet_item_strip = bullet_item.strip()
        p = tf.add_paragraph()
        p.text = bullet_item_strip
        p.level = 1
    image_shape = slide.placeholders[2]
    slide.shapes.add_picture(item.image, image_shape.left, image_shape.top, image_shape.width, image_shape.height)
    return "added text and image page."


@router.get("/tools/ppt/submit_file")
async def submit_file():
    global ppt_file
    file_path = os.path.join(CACHE_DIR, f"{time.time()}.pptx")
    ppt_file.save(file_path)
    return f"submitted. view ppt at {file_path}"
